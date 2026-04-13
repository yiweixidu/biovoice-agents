"""
core/presentation/ppt_generator.py
Google-Slides-quality PPT generator built on python-pptx primitives.

Design language
  - Color palette: navy #1A3557 | teal accent #2E86AB | light gray #F5F6FA
    text #1C1C1E | subtext #6E7B8B
  - Fonts: Calibri (titles 32–40pt bold, body 14–18pt, captions 11pt)
  - Title slide: full-bleed navy header band, white title, teal rule
  - Content slides: left accent bar (4px teal), title top-left, body below
  - Table slides: zebra rows with teal header
  - All slides: 16:9 widescreen (13.33 x 7.5 in)
"""

from __future__ import annotations

import io
import os
import textwrap
from datetime import date
from typing import TYPE_CHECKING, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

if TYPE_CHECKING:
    from biovoice.core.grant_config import GrantOutput

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x35, 0x57)
TEAL   = RGBColor(0x2E, 0x86, 0xAB)
LGRAY  = RGBColor(0xF5, 0xF6, 0xFA)
MGRAY  = RGBColor(0xD0, 0xD3, 0xDA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT   = RGBColor(0x1C, 0x1C, 0x1E)
SUBTEXT= RGBColor(0x6E, 0x7B, 0x8B)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def _rgb(color: RGBColor):
    return color


def _add_rect(slide, l, t, w, h, fill: RGBColor, line: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text: str, l, t, w, h,
               size=Pt(14), bold=False, color=TEXT,
               align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return txb


def _add_slide_number(slide, num: int):
    """Bottom-right page number."""
    _add_text(
        slide, str(num),
        W - Inches(0.6), H - Inches(0.35),
        Inches(0.5), Inches(0.3),
        size=Pt(9), color=SUBTEXT, align=PP_ALIGN.RIGHT,
    )


class PPTGenerator:
    """
    Programmatic slide builder.  No template required.
    All slides are constructed from raw shapes so the design is fully
    controllable without a .pptx template file.
    """

    def __init__(self, template_path: Optional[str] = None):
        self.prs = Presentation()
        self.prs.slide_width  = W
        self.prs.slide_height = H
        self._slide_num = 0

    # ── Internal slide factory ────────────────────────────────────────────────

    def _blank(self):
        """Return a truly blank slide (blank layout, no placeholders)."""
        layout = self.prs.slide_layouts[6]   # blank
        slide  = self.prs.slides.add_slide(layout)
        # Remove any leftover placeholders from the layout
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        self._slide_num += 1
        return slide

    def _header_band(self, slide, height=Inches(1.25)):
        """Full-width navy header band."""
        _add_rect(slide, 0, 0, W, height, NAVY)
        return height

    def _teal_rule(self, slide, top, width=Inches(1.2), height=Pt(4)):
        """Short horizontal teal rule used under section titles."""
        _add_rect(slide, Inches(0.55), top, width, height, TEAL)

    def _slide_bg(self, slide):
        """Light gray slide background."""
        _add_rect(slide, 0, 0, W, H, LGRAY)

    def _accent_bar(self, slide):
        """Left vertical teal accent bar for content slides."""
        _add_rect(slide, 0, 0, Pt(5), H, TEAL)

    # ── Public slide builders ─────────────────────────────────────────────────

    def add_title_slide(self, title: str, subtitle: str = ""):
        slide = self._blank()
        # Background
        _add_rect(slide, 0, 0, W, H, WHITE)
        # Top navy band (40% of height)
        band_h = Inches(3.1)
        _add_rect(slide, 0, 0, W, band_h, NAVY)
        # Teal rule below band
        _add_rect(slide, 0, band_h, W, Pt(5), TEAL)
        # Title in band
        _add_text(
            slide, title,
            Inches(0.6), Inches(0.55), Inches(12.0), Inches(2.2),
            size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.LEFT,
        )
        # Subtitle below band
        if subtitle:
            _add_text(
                slide, subtitle,
                Inches(0.6), band_h + Inches(0.35), Inches(11.5), Inches(2.5),
                size=Pt(18), color=SUBTEXT, align=PP_ALIGN.LEFT,
            )
        # Bottom footer strip
        _add_rect(slide, 0, H - Inches(0.45), W, Inches(0.45), LGRAY)
        _add_text(
            slide, f"Generated {date.today().strftime('%B %d, %Y')}",
            Inches(0.6), H - Inches(0.4), Inches(6), Inches(0.38),
            size=Pt(10), color=SUBTEXT,
        )
        _add_slide_number(slide, self._slide_num)
        return slide

    def add_section_divider(self, section_name: str):
        """Full-bleed teal section divider slide."""
        slide = self._blank()
        _add_rect(slide, 0, 0, W, H, TEAL)
        _add_text(
            slide, section_name,
            Inches(1.0), Inches(2.8), Inches(11.0), Inches(1.8),
            size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )
        _add_slide_number(slide, self._slide_num)
        return slide

    def add_content_slide(self, title: str, bullets: List[str],
                          subtitle: str = ""):
        slide = self._blank()
        _add_rect(slide, 0, 0, W, H, WHITE)
        self._accent_bar(slide)
        # Header band
        band_h = Inches(1.3)
        _add_rect(slide, Pt(5), 0, W - Pt(5), band_h, NAVY)
        # Title
        _add_text(
            slide, title,
            Inches(0.55), Inches(0.12), Inches(11.8), Inches(1.0),
            size=Pt(28), bold=True, color=WHITE,
        )
        if subtitle:
            _add_text(
                slide, subtitle,
                Inches(0.55), Inches(0.95), Inches(11.8), Inches(0.35),
                size=Pt(12), color=RGBColor(0xB0,0xC4,0xD8),
            )
        # Body
        body_top = band_h + Inches(0.25)
        body_h   = H - body_top - Inches(0.5)
        body_txb = slide.shapes.add_textbox(
            Inches(0.65), body_top, Inches(12.3), body_h
        )
        tf = body_txb.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"\u2022  {bullet}"
            run.font.size  = Pt(15)
            run.font.color.rgb = TEXT
            run.font.name  = "Calibri"
        # Footer
        _add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), LGRAY)
        _add_slide_number(slide, self._slide_num)
        return slide

    def add_prose_slide(self, title: str, body: str, citation_note: str = ""):
        """Content slide with long prose text (wraps automatically)."""
        slide = self._blank()
        _add_rect(slide, 0, 0, W, H, WHITE)
        self._accent_bar(slide)
        band_h = Inches(1.15)
        _add_rect(slide, Pt(5), 0, W - Pt(5), band_h, NAVY)
        _add_text(
            slide, title,
            Inches(0.55), Inches(0.15), Inches(11.8), Inches(0.95),
            size=Pt(26), bold=True, color=WHITE,
        )
        # Wrap body text to ~130 chars per line
        body_top = band_h + Inches(0.2)
        footer_h = Inches(0.45) if citation_note else Inches(0.35)
        body_h   = H - body_top - footer_h
        body_txb = slide.shapes.add_textbox(
            Inches(0.65), body_top, Inches(12.3), body_h,
        )
        tf = body_txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        # Trim to what fits (~900 chars for 15pt text)
        display = body[:920] + (" …" if len(body) > 920 else "")
        run.text = display
        run.font.size  = Pt(14)
        run.font.color.rgb = TEXT
        run.font.name  = "Calibri"
        # Citation note
        if citation_note:
            _add_text(
                slide, citation_note,
                Inches(0.65), H - footer_h, Inches(10), Inches(0.38),
                size=Pt(10), color=SUBTEXT, italic=True,
            )
        _add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), LGRAY)
        _add_slide_number(slide, self._slide_num)
        return slide

    def add_two_column_slide(self, title: str, left_title: str, left_body: str,
                              right_title: str, right_body: str):
        slide = self._blank()
        _add_rect(slide, 0, 0, W, H, WHITE)
        self._accent_bar(slide)
        band_h = Inches(1.15)
        _add_rect(slide, Pt(5), 0, W - Pt(5), band_h, NAVY)
        _add_text(slide, title,
                  Inches(0.55), Inches(0.15), Inches(11.8), Inches(0.95),
                  size=Pt(26), bold=True, color=WHITE)
        # Divider
        mid = W / 2
        _add_rect(slide, mid - Pt(0.5), band_h + Inches(0.1),
                  Pt(1), H - band_h - Inches(0.5), MGRAY)
        col_w = mid - Inches(0.85)
        for lx, lt, lb in [
            (Inches(0.65), left_title,  left_body),
            (mid + Inches(0.2), right_title, right_body),
        ]:
            _add_text(slide, lt, lx, band_h + Inches(0.2), col_w, Inches(0.4),
                      size=Pt(14), bold=True, color=TEAL)
            _add_text(slide, lb, lx, band_h + Inches(0.65), col_w,
                      H - band_h - Inches(1.1),
                      size=Pt(13), color=TEXT)
        _add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), LGRAY)
        _add_slide_number(slide, self._slide_num)
        return slide

    def add_table_slide(self, title: str, headers: List[str],
                        rows: List[List[str]]):
        slide = self._blank()
        _add_rect(slide, 0, 0, W, H, WHITE)
        self._accent_bar(slide)
        band_h = Inches(1.15)
        _add_rect(slide, Pt(5), 0, W - Pt(5), band_h, NAVY)
        _add_text(slide, title,
                  Inches(0.55), Inches(0.15), Inches(11.8), Inches(0.95),
                  size=Pt(26), bold=True, color=WHITE)

        rows_count = min(len(rows), 14)
        cols_count = len(headers)
        row_h = min(Inches(0.38), (H - band_h - Inches(0.6)) / (rows_count + 1))
        tbl_h = row_h * (rows_count + 1)
        table = slide.shapes.add_table(
            rows_count + 1, cols_count,
            Inches(0.55), band_h + Inches(0.15),
            W - Inches(0.7), tbl_h,
        ).table

        def _cell_style(cell, text, bg: RGBColor, fg: RGBColor,
                        bold=False, size=Pt(11)):
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size  = size
                    run.font.bold  = bold
                    run.font.color.rgb = fg
                    run.font.name  = "Calibri"

        for col, hdr in enumerate(headers):
            _cell_style(table.cell(0, col), hdr, NAVY, WHITE, bold=True, size=Pt(11))

        for i, row in enumerate(rows[:rows_count]):
            bg = WHITE if i % 2 == 0 else LGRAY
            for j, val in enumerate(row):
                _cell_style(table.cell(i + 1, j), str(val), bg, TEXT, size=Pt(10))

        _add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), LGRAY)
        _add_slide_number(slide, self._slide_num)
        return slide

    def add_chart_slide(self, title: str, img_bytes, caption: str = ""):
        """
        Embed any matplotlib chart (as BytesIO PNG) into a slide.
        The chart fills the content area below the navy header band.
        """
        slide = self._blank()
        _add_rect(slide, 0, 0, W, H, WHITE)
        self._accent_bar(slide)
        band_h = Inches(1.15)
        _add_rect(slide, Pt(5), 0, W - Pt(5), band_h, NAVY)
        _add_text(slide, title,
                  Inches(0.55), Inches(0.15), Inches(11.8), Inches(0.95),
                  size=Pt(26), bold=True, color=WHITE)
        # Centre the chart image in the remaining content area
        content_top  = band_h + Inches(0.1)
        content_h    = H - content_top - Inches(0.5)
        slide.shapes.add_picture(img_bytes,
                                 Inches(0.6), content_top,
                                 height=content_h)
        if caption:
            _add_text(slide, caption,
                      Inches(0.65), H - Inches(0.45), Inches(12), Inches(0.38),
                      size=Pt(9), color=SUBTEXT, italic=True)
        _add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), LGRAY)
        _add_slide_number(slide, self._slide_num)
        return slide

    def add_neutralization_heatmap(self, data: Dict[str, List[float]],
                                   title: str = "Neutralization Breadth",
                                   strain_labels: Optional[List[str]] = None):
        from .visualizer import create_neutralization_heatmap
        img_bytes = create_neutralization_heatmap(data, title, strain_labels)
        return self.add_chart_slide(title, img_bytes)

    def save(self, path: str):
        os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
        self.prs.save(path)

    # ── Grant mode ────────────────────────────────────────────────────────────

    def add_grant_section_slide(self, section_title: str, text: str,
                                citation_count: int = 0):
        note = f"{citation_count} source(s) cited" if citation_count else ""
        self.add_prose_slide(section_title, text, note)

    def add_grant_key_findings_slide(self, papers):
        if not papers:
            return
        headers = ["#", "Authors", "Year", "Title", "PMID"]
        rows = [
            [
                str(i + 1),
                (c.authors or "")[:28],
                c.year or "",
                (c.title or "")[:58],
                c.pmid or "",
            ]
            for i, c in enumerate(papers[:14])
        ]
        self.add_table_slide("Key Literature", headers, rows)

    def add_grant_references_slide(self, citations):
        refs = []
        for i, c in enumerate(citations[:18], 1):
            refs.append(
                f"[{i}] {(c.authors or 'Unknown')[:35]}. "
                f"{(c.title or '')[:55]}. "
                f"{c.journal or ''} {c.year or ''}. PMID:{c.pmid}"
            )
        self.add_content_slide("References", refs)


# ── Standalone render function (used by grant pipeline) ───────────────────────

def render_grant_ppt(grant_output: "GrantOutput", output_path: str) -> str:
    gen = PPTGenerator()

    gen.add_title_slide(
        "Virology Grant Copilot",
        grant_output.research_question[:120],
    )

    section_order = ["specific_aims", "significance", "innovation", "background"]
    for key in section_order:
        gs = grant_output.section(key)
        if gs is None:
            continue
        gen.add_grant_section_slide(gs.title, gs.text, len(gs.citations))

    gen.add_grant_key_findings_slide(grant_output.all_citations)
    gen.add_grant_references_slide(grant_output.all_citations)

    gen.save(output_path)
    print(f"[PPT] Saved: {output_path}")
    return output_path
